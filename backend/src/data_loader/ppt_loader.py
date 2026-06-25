from __future__ import annotations

from pathlib import Path
from typing import Iterator, Optional

from langchain_core.documents import Document
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.table import Table

from src.data_loader.base_loader import BaseDocumentLoader
from src.data_loader.loader_registry import loader_registry


@loader_registry.register(extensions=[".pptx"])
class PptxLoader(BaseDocumentLoader):
    def lazy_load(self) -> Iterator[Document]:
        prs = Presentation(str(self._file_path))

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            texts.append(f"[Slide {slide_num}]")

            for shape in slide.shapes:
                content = self._extract_shape_text(shape)
                if content:
                    texts.append(content)

            combined = "\n".join(texts)
            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            if notes:
                combined += f"\n\n[Speaker Notes]\n{notes}"

            if len(texts) > 1:
                yield Document(
                    page_content=combined,
                    metadata={
                        "slide_number": slide_num,
                        "total_slides": len(prs.slides),
                        "source_file": self._file_path.name,
                        "has_notes": bool(notes),
                    },
                )

    def _extract_shape_text(self, shape: BaseShape) -> str:
        if shape.has_text_frame:
            return shape.text_frame.text.strip()
        if shape.has_table:
            return self._extract_table(shape.table)
        if hasattr(shape, 'shapes'):
            parts = []
            for s in shape.shapes:
                t = self._extract_shape_text(s)
                if t:
                    parts.append(t)
            return "\n".join(parts)
        return ""

    def _extract_table(self, table: Table) -> str:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "[Table]\n" + "\n".join(rows)

    @classmethod
    def supports(cls, file_path: Path) -> bool:
        return file_path.suffix.lower() in (".pptx", ".ppt")